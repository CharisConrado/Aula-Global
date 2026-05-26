"""
Aula Global — Seed Currículo Primaria
Extrae contenido de 120 PPTXs, los sube a Supabase Storage
y genera actividades evaluativas por materia.

Actividades por materia:
  Arte              → Dibujo (canvas)
  Ciencias Sociales → Subir video
  Inglés            → Selección múltiple (quiz)
  Lenguaje/Español  → Emparejar conceptos
  Matemáticas       → Ejercicios + subir archivo
  Ciencias Naturales→ Quiz

Uso:
    cd backend
    pip install python-pptx supabase psycopg2-binary python-dotenv
    python seed_curriculum.py
"""
import os
import json
import sys
import random
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

SUPABASE_URL         = os.getenv("SUPABASE_URL", "")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY", "")
DATABASE_URL         = os.getenv("DATABASE_URL", "")

TEMAS_DIR = Path(r"D:\Users\TRABAJO\Downloads\Curriculo_Primaria_120_Temas_1\temas")
BUCKET    = "presentations"

random.seed(7)  # reproducible

# ── Mapeos ────────────────────────────────────────────────────────────────────

GRADE_FOLDER_TO_LEVEL = {
    "1-primero-de-primaria": 1,
    "2-segundo-de-primaria": 2,
    "3-tercero-de-primaria": 3,
    "4-cuarto-de-primaria":  4,
    "5-quinto-de-primaria":  5,
}

SUBJECT_KEY_MAP = {
    "arte":               ("Arte",               "🎨", "#FF6B6B"),
    "ciencias-naturales": ("Ciencias Naturales", "🌿", "#43A047"),
    "ciencias-sociales":  ("Ciencias Sociales",  "🌍", "#FF9800"),
    "ingles":             ("Inglés",             "🇺🇸", "#2196F3"),
    "lenguaje":           ("Lenguaje",           "📚", "#9C27B0"),
    "matematicas":        ("Matemáticas",        "🔢", "#E53935"),
}

# ── Qué tipo de actividad usa cada materia ─────────────────────────────────────
# tipo_evaluacion → lo que el frontend usa para saber qué UI mostrar
# db_type_name    → nombre en la tabla type_activity de la DB
SUBJECT_ACTIVITY = {
    "arte":               {"tipo_eval": "dibujo",           "db_type": "dibujo"},
    "ciencias-sociales":  {"tipo_eval": "video_upload",     "db_type": "video"},
    "ingles":             {"tipo_eval": "quiz",              "db_type": "quiz"},
    "lenguaje":           {"tipo_eval": "emparejar",        "db_type": "asociar"},
    "matematicas":        {"tipo_eval": "ejercicio_archivo","db_type": "ejercicio"},
    "ciencias-naturales": {"tipo_eval": "quiz",              "db_type": "quiz"},
}

DIFFICULTY_ROTATION = ["facil", "medio", "dificil"]

STOP_WORDS = {
    "y", "e", "o", "u", "a", "de", "del", "la", "el", "los",
    "las", "un", "una", "con", "en", "al",
}

# ── Texto ─────────────────────────────────────────────────────────────────────

def slug_to_title(slug: str) -> str:
    words = slug.replace("-", " ").split()
    return " ".join(
        w.capitalize() if i == 0 or w not in STOP_WORDS else w
        for i, w in enumerate(words)
    )

def parse_filename(name: str):
    stem  = Path(name).stem
    parts = stem.split("_", 2)
    if len(parts) < 3:
        return None, stem
    return parts[1], slug_to_title(parts[2])

# ── Extracción PPTX ───────────────────────────────────────────────────────────

TITLE_PH_TYPES = {1, 3}  # PP_PLACEHOLDER.TITLE=1, CENTER_TITLE=3

def extract_slides(pptx_path: Path):
    """Extrae el texto de cada diapositiva. Devuelve (slides_list, full_text)."""
    try:
        from pptx import Presentation
    except ImportError:
        return [], ""

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f"    ! No se pudo leer PPTX: {e}")
        return [], ""

    slides_data = []
    all_parts   = []

    for slide_idx, slide in enumerate(prs.slides):
        titulo = ""
        puntos = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            is_title = False
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.type in TITLE_PH_TYPES:
                    is_title = True
            except Exception:
                pass

            if is_title:
                titulo = text
            else:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line and len(line) > 2:
                        puntos.append(line[:220])

        if titulo or puntos:
            slides_data.append({
                "num":    slide_idx + 1,
                "titulo": titulo,
                "puntos": puntos[:8],
            })
            if titulo:
                all_parts.append(titulo)
            all_parts.extend(puntos[:3])

    return slides_data, " | ".join(all_parts)

# ── Generadores de contenido por materia ─────────────────────────────────────

def _bullets(slides):
    out = []
    for s in slides:
        out.extend(s["puntos"])
    return out

def _titles(slides):
    return [s["titulo"] for s in slides if s["titulo"]]


def gen_arte(title, slides):
    topic = _titles(slides)[0] if _titles(slides) else title
    return {
        "tipo_evaluacion": "dibujo",
        "instruccion": (
            f"🎨 Dibuja lo que aprendiste sobre '{topic}'. "
            "Usa todos los colores que quieras y sé creativo."
        ),
    }


def gen_sociales(title, slides):
    key   = _titles(slides)[0] if _titles(slides) else title
    puntos = _bullets(slides)[:3]
    return {
        "tipo_evaluacion": "video_upload",
        "instruccion": (
            f"Graba un video corto (1-2 minutos) explicando con tus propias palabras "
            f"qué es '{key}' y por qué es importante para la sociedad."
        ),
        "puntos_clave": puntos,
    }


def gen_ingles(title, slides):
    bullets  = _bullets(slides)
    preguntas = []

    for s in slides:
        if not s["puntos"]:
            continue
        correct = s["puntos"][0]
        pool    = [b for b in bullets if b != correct]
        random.shuffle(pool)
        distractors = pool[:3]
        while len(distractors) < 3:
            distractors.append("None of the above")

        opts = [correct] + distractors
        random.shuffle(opts)
        correct_idx = opts.index(correct)
        topic = s["titulo"] or title

        preguntas.append({
            "pregunta":           f"Which statement is correct about '{topic}'?",
            "opciones":           opts,
            "respuesta_correcta": correct_idx,
            "pista":              f"Review slide {s['num']}",
        })
        if len(preguntas) >= 5:
            break

    if not preguntas:
        preguntas = [{
            "pregunta":           f"What is the main topic of this presentation?",
            "opciones":           [title, "A different topic", "Ancient history", "A story"],
            "respuesta_correcta": 0,
        }]

    return {
        "tipo_evaluacion": "quiz",
        "preguntas":        preguntas,
    }


def gen_lenguaje(title, slides):
    pares = []
    for s in slides:
        if s["titulo"] and s["puntos"]:
            pares.append({"a": s["titulo"], "b": s["puntos"][0]})
        if len(pares) >= 6:
            break

    if not pares:
        bullets = _bullets(slides)
        titles  = _titles(slides)
        for i in range(min(len(titles), len(bullets), 4)):
            pares.append({"a": titles[i], "b": bullets[i]})

    if not pares:
        pares = [
            {"a": title,        "b": "Tema principal de la lección"},
            {"a": "Definición", "b": f"Concepto clave de {title}"},
        ]

    return {
        "tipo_evaluacion": "emparejar",
        "pares":            pares,
    }


def gen_matematicas(title, slides):
    bullets    = _bullets(slides)
    ejercicios = []

    for bullet in bullets:
        ejercicios.append(bullet)
        if len(ejercicios) >= 5:
            break

    # Si no hay suficientes bullets, genera ejercicios genéricos del tema
    if len(ejercicios) < 3:
        key = _titles(slides)[0] if _titles(slides) else title
        ejercicios = [
            f"Resuelve 3 problemas sobre {key} usando lo aprendido en la presentación.",
            f"Explica con un ejemplo propio el concepto de {title}.",
            f"Dibuja o escribe los pasos para resolver un ejercicio de {key}.",
        ]

    return {
        "tipo_evaluacion": "ejercicio_archivo",
        "ejercicios":       ejercicios[:5],
        "instruccion": (
            "Resuelve los ejercicios en tu cuaderno o en una hoja, "
            "luego toma una foto o sube el archivo PDF con tus respuestas."
        ),
    }


def gen_naturales(title, slides):
    # Igual que inglés pero en español
    bullets   = _bullets(slides)
    preguntas = []

    for s in slides:
        if not s["puntos"]:
            continue
        correct = s["puntos"][0]
        pool    = [b for b in bullets if b != correct]
        random.shuffle(pool)
        distractors = pool[:3]
        while len(distractors) < 3:
            distractors.append("No se menciona en la presentación")

        opts = [correct] + distractors
        random.shuffle(opts)

        preguntas.append({
            "pregunta":           f"¿Qué es correcto sobre '{s['titulo'] or title}'?",
            "opciones":           opts,
            "respuesta_correcta": opts.index(correct),
            "pista":              f"Revisa la diapositiva {s['num']}",
        })
        if len(preguntas) >= 5:
            break

    if not preguntas:
        preguntas = [{
            "pregunta":           f"¿De qué trata esta presentación de Ciencias Naturales?",
            "opciones":           [title, "Historia del mundo", "Arte y colores", "Geografía"],
            "respuesta_correcta": 0,
        }]

    return {
        "tipo_evaluacion": "quiz",
        "preguntas":        preguntas,
    }


CONTENT_GENERATORS = {
    "arte":               gen_arte,
    "ciencias-sociales":  gen_sociales,
    "ingles":             gen_ingles,
    "lenguaje":           gen_lenguaje,
    "matematicas":        gen_matematicas,
    "ciencias-naturales": gen_naturales,
}


def generate_content(subject_key: str, title: str, slides: list, public_url: str) -> dict:
    generator = CONTENT_GENERATORS.get(subject_key, gen_naturales)
    content   = generator(title, slides)
    # Siempre incluir slides y URL para el visor
    content["slides"]           = slides
    content["presentacion_url"] = public_url
    return content

# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    try:
        from supabase import create_client
        import psycopg2
    except ImportError as e:
        print(f"Falta dependencia: {e}")
        print("Instala con:  pip install supabase psycopg2-binary python-dotenv python-pptx")
        sys.exit(1)

    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("Falta python-pptx.  pip install python-pptx")
        sys.exit(1)

    print("=" * 65)
    print("  Aula Global — Seed Currículo Primaria")
    print("  Arte→Dibujo | Sociales→Video | Inglés→Quiz")
    print("  Lenguaje→Emparejar | Mates→Archivo | Naturales→Quiz")
    print("=" * 65)

    # ── Supabase Storage ──────────────────────────────────────────────────────
    supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
    try:
        supabase.storage.create_bucket(BUCKET, options={"public": True})
        print(f"✓ Bucket '{BUCKET}' creado")
    except Exception as e:
        msg = str(e).lower()
        if "already exists" in msg or "duplicate" in msg or "409" in msg:
            print(f"• Bucket '{BUCKET}' ya existe — OK")
        else:
            print(f"! Error creando bucket: {e}")

    # ── Base de datos ─────────────────────────────────────────────────────────
    conn = psycopg2.connect(DATABASE_URL)
    cur  = conn.cursor()

    cur.execute("SELECT id_degree, grade_name, level FROM degree ORDER BY level")
    rows = cur.fetchall()
    level_to_degree_id = {r[2]: str(r[0]) for r in rows}
    print(f"\nGrados en DB : {[r[1] for r in rows]}")

    if not level_to_degree_id:
        print("¡No hay grados en la DB!")
        sys.exit(1)

    cur.execute("SELECT id_type_activity, LOWER(name) FROM type_activity")
    type_map = {r[1]: str(r[0]) for r in cur.fetchall()}
    print(f"Tipos en DB  : {list(type_map.keys())}")

    # ── Procesar grados ───────────────────────────────────────────────────────
    total_uploaded = 0
    total_created  = 0
    diff_idx       = 0

    for grade_folder, grade_level in GRADE_FOLDER_TO_LEVEL.items():
        grade_path = TEMAS_DIR / grade_folder
        if not grade_path.exists():
            print(f"\n[!] Carpeta no encontrada: {grade_path}")
            continue

        degree_id = level_to_degree_id.get(grade_level)
        if not degree_id:
            print(f"\n[!] Grado nivel {grade_level} no en DB")
            continue

        print(f"\n── Grado {grade_level} {'─'*45}")

        cur.execute(
            "SELECT id_subject, subject_name FROM subject WHERE id_degree = %s::uuid",
            (degree_id,)
        )
        subject_cache = {r[1]: str(r[0]) for r in cur.fetchall()}

        for pptx_file in sorted(grade_path.glob("*.pptx")):
            subject_key, title = parse_filename(pptx_file.name)

            if subject_key not in SUBJECT_KEY_MAP:
                print(f"  [!] Subject key desconocido '{subject_key}' — {pptx_file.name}")
                continue

            subject_name, icon, color = SUBJECT_KEY_MAP[subject_key]

            # Crear materia si no existe
            if subject_name not in subject_cache:
                cur.execute("""
                    INSERT INTO subject (id_degree, subject_name, description, icon, color, is_active)
                    VALUES (%s::uuid, %s, %s, %s, %s, true)
                    RETURNING id_subject
                """, (degree_id, subject_name, f"Materia de {subject_name}", icon, color))
                subject_id = str(cur.fetchone()[0])
                conn.commit()
                subject_cache[subject_name] = subject_id
                print(f"  ✓ Materia creada: {subject_name}")
            else:
                subject_id = subject_cache[subject_name]

            # ── Extraer slides del PPTX ───────────────────────────────────────
            slides, _ = extract_slides(pptx_file)
            print(f"  → {pptx_file.name}  [{len(slides)} diap.]")

            # ── Subir PPTX a Supabase Storage ─────────────────────────────────
            storage_path = f"{grade_folder}/{pptx_file.name}"
            public_url   = f"{SUPABASE_URL}/storage/v1/object/public/{BUCKET}/{storage_path}"

            try:
                with open(pptx_file, "rb") as f:
                    file_bytes = f.read()

                supabase.storage.from_(BUCKET).upload(
                    storage_path,
                    file_bytes,
                    {
                        "content-type": (
                            "application/vnd.openxmlformats-officedocument"
                            ".presentationml.presentation"
                        ),
                        "upsert": "true",
                    },
                )
                total_uploaded += 1
                print(f"     ↑ Subido a Storage")
            except Exception as e:
                msg = str(e).lower()
                if "already exists" in msg or "duplicate" in msg or "409" in msg:
                    print(f"     • Ya en Storage")
                else:
                    print(f"     [!] Error subiendo: {e}")

            # ── Tipo de actividad según materia ────────────────────────────────
            subject_cfg = SUBJECT_ACTIVITY.get(subject_key, {"tipo_eval": "quiz", "db_type": "quiz"})
            db_type_name = subject_cfg["db_type"]

            # Buscar el id del tipo en la DB (coincidencia parcial)
            type_id = type_map.get(db_type_name)
            if not type_id:
                for k, v in type_map.items():
                    if db_type_name in k or k in db_type_name:
                        type_id = v
                        break
            if not type_id:
                print(f"     [!] Tipo '{db_type_name}' no en DB — omitiendo")
                continue

            # ── Generar contenido evaluativo por materia ───────────────────────
            activity_content = generate_content(subject_key, title, slides, public_url)
            content_str      = json.dumps(activity_content, ensure_ascii=False)

            difficulty = DIFFICULTY_ROTATION[diff_idx % len(DIFFICULTY_ROTATION)]
            diff_idx  += 1

            # ── Insertar actividad en la DB ────────────────────────────────────
            tipo_eval = subject_cfg["tipo_eval"]
            try:
                cur.execute("""
                    INSERT INTO activity
                        (id_subject, id_type_activity, title, description,
                         difficulty_level, content, estimated_minutes, publication_status)
                    VALUES
                        (%s::uuid, %s::uuid, %s, %s,
                         %s, %s::jsonb, 20, 'publicado')
                """, (
                    subject_id, type_id,
                    title, f"Aprende sobre {title}",
                    difficulty, content_str,
                ))
                conn.commit()
                total_created += 1
                print(f"     ✓ [{tipo_eval}/{difficulty}] {title}")
            except Exception as e:
                conn.rollback()
                print(f"     [!] Error creando actividad '{title}': {e}")

    cur.close()
    conn.close()

    print("\n" + "=" * 65)
    print(f"  ✓ Archivos subidos    : {total_uploaded}")
    print(f"  ✓ Actividades creadas : {total_created}")
    print("=" * 65)
    print("\n  Arte→Dibujo ✏️  |  Sociales→Video 🎥  |  Inglés→Quiz 🧠")
    print("  Lenguaje→Emparejar 🔗  |  Mates→Archivo 📁  |  Naturales→Quiz 🌿")

if __name__ == "__main__":
    main()
